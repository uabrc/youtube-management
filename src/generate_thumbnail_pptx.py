import collections  # ! supports pptx
import collections.abc  # ! support pptx
import textwrap
from pathlib import PurePath
from typing import Dict

import pandas as pd
import pptx

# FIELDS
TITLE = "title"
SUBTITLE = "subtitle"
PART = "part"
CATEGORY = "category"
INDEX = "index"
DATE = "date"
LAYOUT_INDEX = "layout_index"

# LAYOUT
TITLE_PLACEHOLDER_INDEX = 0
TITLE_FONT_SIZE = pptx.util.Pt(48)  # type: ignore
SUBTITLE_PLACEHOLDER_INDEX = 1
SUBTITLE_FONT_SIZE = pptx.util.Pt(24)  # type: ignore


class ContentBuilder:
    def __init__(self, organization_name: str):
        self._organization_name: str = organization_name

    def build_content(self, row: pd.Series) -> Dict[str, str]:
        return {
            "title": self._build_title_content(row=row),
            "subtitle": self._build_subtitle_content(row=row),
        }

    def _build_title_content(self, row: pd.Series) -> str:
        part = row[PART]
        part_content = self._build_part_content(part=part)
        title_content = textwrap.dedent(
            f"""
            {row[TITLE]}{part_content}
            """
        ).strip()
        return title_content

    def _build_part_content(self, part: pd.Series) -> str:
        if not pd.isna(part):
            part_value = int(str(part))
            part_content = f" (Part {part_value:d})"
        else:
            part_content = ""
        return part_content

    def _build_subtitle_content(self, row: pd.Series) -> str:
        category_content = self._build_category_content(row[CATEGORY])
        index_content = self._build_index_content(row[INDEX])

        date = row[DATE].strftime(r"%Y-%m-%d")

        subtitle_content = textwrap.dedent(
            f"""
            {category_content}{index_content}
            {self._organization_name}
            {date}
            """
        ).strip()
        return subtitle_content

    def _build_category_content(self, category: pd.Series) -> str:
        if not pd.isna(category):
            category_content = str(category)
        else:
            category_content = ""
        return category_content

    def _build_index_content(self, index: pd.Series) -> str:
        if not pd.isna(index):
            index_value = int(str(index))
            index_content = f" #{index_value:d}"
        else:
            index_content = ""
        return index_content


def preprocess(df: pd.DataFrame) -> pd.DataFrame:
    df[PART] = pd.to_numeric(df[PART])
    df[INDEX] = pd.to_numeric(df[INDEX])
    df[DATE] = pd.to_datetime(df[DATE])
    return df


def add_content_to_slide(content: Dict[str, str], slide: pptx.Slide):
    # TITLE
    title_content = content[TITLE]
    title_placeholder = slide.placeholders[TITLE_PLACEHOLDER_INDEX]
    title_placeholder.text = title_content
    for paragraph in title_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = TITLE_FONT_SIZE

    # SUBTITLE
    subtitle_content = content[SUBTITLE]
    subtitle_placeholder = slide.placeholders[SUBTITLE_PLACEHOLDER_INDEX]
    subtitle_placeholder.text = subtitle_content
    for paragraph in subtitle_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = SUBTITLE_FONT_SIZE


def main():
    # ORGANIZATION CONFIGURATION
    ORGANIZATION_NAME = "UAB IT Research Computing"

    # FILE PATHS
    INPUT_CSV_NAME = PurePath("res") / "thumbnails.csv"
    INPUT_PPTX_TEMPLATE_NAME = PurePath("res") / "youtube-thumbnail-template.pptx"
    OUTPUT_PPTX_NAME = PurePath("thumbnails.pptx")

    # PRESENTATION CONFIGURATION
    PPTX_16_9_EMU = (12192000, 6868000)

    presentation = pptx.Presentation(INPUT_PPTX_TEMPLATE_NAME)
    presentation.slide_width = PPTX_16_9_EMU[0]
    presentation.slide_height = PPTX_16_9_EMU[1]

    df = pd.read_csv(INPUT_CSV_NAME)
    df = preprocess(df=df)

    content_builder = ContentBuilder(organization_name=ORGANIZATION_NAME)
    for _, row in df.iterrows():
        layout_index = row[LAYOUT_INDEX]
        slide_layout = presentation.slide_layouts[layout_index]
        slide = presentation.slides.add_slide(slide_layout)
        content = content_builder.build_content(row=row)
        add_content_to_slide(content=content, slide=slide)

    presentation.save(OUTPUT_PPTX_NAME)


if __name__ == "__main__":
    main()
